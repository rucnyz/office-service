"""
Lightweight Office API service for Windows VM.

Provides REST endpoints to create and manipulate Word, Excel, and PowerPoint
files inside the VM. Designed to be called by the injection MCP server on the
host for red-teaming environment setup.

Usage:
    uv run main.py [--port PORT]
    # or
    python main.py [--port PORT]
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

import uvicorn
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field

app = FastAPI(
    title="Office Service API",
    description="Lightweight API for creating/reading Office documents in the VM",
    version="0.2.0",
)


# ── Request / Response models ────────────────────────────────


class ApiResponse(BaseModel):
    status: str = "success"
    message: str = ""
    data: Optional[dict] = None


# -- Word models --

class CreateWordRequest(BaseModel):
    file_path: str
    title: Optional[str] = None
    paragraphs: List[str] = Field(default_factory=list)


class AddWordContentRequest(BaseModel):
    file_path: str
    paragraphs: List[str] = Field(default_factory=list)
    headings: Optional[List[dict]] = None


class SearchReplaceWordRequest(BaseModel):
    file_path: str
    search: str
    replace: str


class ReadWordRequest(BaseModel):
    file_path: str


class AddWordTableRequest(BaseModel):
    file_path: str
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    style: Optional[str] = None


class FormatWordTextRequest(BaseModel):
    file_path: str
    paragraph_index: int
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color: Optional[str] = None
    font_size: Optional[int] = None
    font_name: Optional[str] = None
    hidden: Optional[bool] = None


class AddWordHyperlinkRequest(BaseModel):
    file_path: str
    text: str
    url: str


class DeleteWordParagraphRequest(BaseModel):
    file_path: str
    paragraph_index: int


class WordPageBreakRequest(BaseModel):
    file_path: str


class AddWordHeaderFooterRequest(BaseModel):
    file_path: str
    header_text: Optional[str] = None
    footer_text: Optional[str] = None


# -- Excel models --

class SheetData(BaseModel):
    name: str = "Sheet1"
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    hidden: bool = False


class CreateExcelRequest(BaseModel):
    file_path: str
    sheets: List[SheetData] = Field(default_factory=list)


class WriteExcelDataRequest(BaseModel):
    file_path: str
    sheet_name: str = "Sheet1"
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    hidden: bool = False


class ReadExcelRequest(BaseModel):
    file_path: str
    sheet_name: Optional[str] = None


class ApplyFormulaRequest(BaseModel):
    file_path: str
    sheet_name: str
    cell: str
    formula: str


class FormatRangeRequest(BaseModel):
    file_path: str
    sheet_name: str
    start_cell: str
    end_cell: Optional[str] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    font_size: Optional[int] = None
    font_color: Optional[str] = None
    fill_color: Optional[str] = None
    number_format: Optional[str] = None
    hidden: Optional[bool] = None


class MergeCellsRequest(BaseModel):
    file_path: str
    sheet_name: str
    start_cell: str
    end_cell: str


class SheetOpRequest(BaseModel):
    file_path: str
    sheet_name: str
    new_name: Optional[str] = None


class InsertRowsColsRequest(BaseModel):
    file_path: str
    sheet_name: str
    index: int
    count: int = 1


# -- PowerPoint models --

class SlideData(BaseModel):
    title: str = ""
    content: str = ""
    notes: str = ""


class CreatePptxRequest(BaseModel):
    file_path: str
    slides: List[SlideData] = Field(default_factory=list)


class AddSlideRequest(BaseModel):
    file_path: str
    title: str = ""
    content: str = ""
    notes: str = ""


class ReadPptxRequest(BaseModel):
    file_path: str


class AddPptxTableRequest(BaseModel):
    file_path: str
    slide_index: int
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    left: float = 1.0
    top: float = 2.0
    width: float = 8.0
    height: float = 3.0


class UpdateSlideContentRequest(BaseModel):
    file_path: str
    slide_index: int
    title: Optional[str] = None
    content: Optional[str] = None
    notes: Optional[str] = None


class DeleteSlideRequest(BaseModel):
    file_path: str
    slide_index: int


class DuplicateSlideRequest(BaseModel):
    file_path: str
    slide_index: int


class SetSlideNotesRequest(BaseModel):
    file_path: str
    slide_index: int
    notes: str


class GetSlideInfoRequest(BaseModel):
    file_path: str
    slide_index: int


class AddPptxShapeRequest(BaseModel):
    file_path: str
    slide_index: int
    shape_type: str = "rectangle"
    left: float = 1.0
    top: float = 1.0
    width: float = 3.0
    height: float = 1.0
    text: str = ""
    fill_color: Optional[str] = None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  Word endpoints
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@app.post("/word/create", response_model=ApiResponse, tags=["Word"])
async def word_create(req: CreateWordRequest):
    """Create a new Word document with optional title and paragraphs."""
    from docx import Document

    try:
        doc = Document()
        if req.title:
            doc.add_heading(req.title, level=1)
        for para in req.paragraphs:
            doc.add_paragraph(para)
        Path(req.file_path).parent.mkdir(parents=True, exist_ok=True)
        doc.save(req.file_path)
        return ApiResponse(message=f"Created {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/add_content", response_model=ApiResponse, tags=["Word"])
async def word_add_content(req: AddWordContentRequest):
    """Add paragraphs and/or headings to an existing Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        if req.headings:
            for h in req.headings:
                doc.add_heading(h.get("text", ""), level=h.get("level", 1))
        for para in req.paragraphs:
            doc.add_paragraph(para)
        doc.save(req.file_path)
        return ApiResponse(message=f"Content added to {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/search_replace", response_model=ApiResponse, tags=["Word"])
async def word_search_replace(req: SearchReplaceWordRequest):
    """Find and replace text in a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        count = 0
        for para in doc.paragraphs:
            if req.search in para.text:
                for run in para.runs:
                    if req.search in run.text:
                        run.text = run.text.replace(req.search, req.replace)
                        count += 1
        doc.save(req.file_path)
        return ApiResponse(message=f"Replaced {count} occurrence(s)")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/read", response_model=ApiResponse, tags=["Word"])
async def word_read(req: ReadWordRequest):
    """Read all text from a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                table_data.append([cell.text for cell in row.cells])
            tables.append(table_data)
        return ApiResponse(
            message="OK",
            data={"paragraphs": paragraphs, "tables": tables},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/add_table", response_model=ApiResponse, tags=["Word"])
async def word_add_table(req: AddWordTableRequest):
    """Add a table to a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        num_cols = len(req.headers) if req.headers else (len(req.rows[0]) if req.rows else 1)
        num_rows = (1 if req.headers else 0) + len(req.rows)
        table = doc.add_table(rows=num_rows, cols=num_cols)
        if req.style:
            table.style = req.style
        row_offset = 0
        if req.headers:
            for j, h in enumerate(req.headers):
                table.cell(0, j).text = h
            row_offset = 1
        for i, row_data in enumerate(req.rows):
            for j, val in enumerate(row_data):
                if j < num_cols:
                    table.cell(i + row_offset, j).text = str(val)
        doc.save(req.file_path)
        return ApiResponse(message=f"Table added ({num_rows}x{num_cols})")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/format_text", response_model=ApiResponse, tags=["Word"])
async def word_format_text(req: FormatWordTextRequest):
    """Format text in a specific paragraph (bold, italic, color, font, hidden)."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    try:
        doc = Document(req.file_path)
        if req.paragraph_index >= len(doc.paragraphs):
            raise HTTPException(status_code=400, detail=f"Paragraph index {req.paragraph_index} out of range")
        para = doc.paragraphs[req.paragraph_index]
        for run in para.runs:
            if req.bold is not None:
                run.font.bold = req.bold
            if req.italic is not None:
                run.font.italic = req.italic
            if req.underline is not None:
                run.font.underline = req.underline
            if req.font_size is not None:
                run.font.size = Pt(req.font_size)
            if req.font_name is not None:
                run.font.name = req.font_name
            if req.color is not None:
                run.font.color.rgb = RGBColor.from_string(req.color.lstrip("#"))
            if req.hidden is not None:
                run.font.hidden = req.hidden
        doc.save(req.file_path)
        return ApiResponse(message=f"Formatted paragraph {req.paragraph_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/add_hyperlink", response_model=ApiResponse, tags=["Word"])
async def word_add_hyperlink(req: AddWordHyperlinkRequest):
    """Add a paragraph with a hyperlink to a Word document."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    try:
        doc = Document(req.file_path)
        para = doc.add_paragraph()

        # Create hyperlink element
        part = doc.part
        r_id = part.relate_to(req.url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        rStyle = OxmlElement("w:rStyle")
        rStyle.set(qn("w:val"), "Hyperlink")
        rPr.append(rStyle)
        run.append(rPr)
        text_elem = OxmlElement("w:t")
        text_elem.text = req.text
        run.append(text_elem)
        hyperlink.append(run)
        para._p.append(hyperlink)

        doc.save(req.file_path)
        return ApiResponse(message=f"Hyperlink added: {req.text}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/delete_paragraph", response_model=ApiResponse, tags=["Word"])
async def word_delete_paragraph(req: DeleteWordParagraphRequest):
    """Delete a paragraph by index from a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        if req.paragraph_index >= len(doc.paragraphs):
            raise HTTPException(status_code=400, detail=f"Paragraph index {req.paragraph_index} out of range")
        p = doc.paragraphs[req.paragraph_index]._p
        p.getparent().remove(p)
        doc.save(req.file_path)
        return ApiResponse(message=f"Deleted paragraph {req.paragraph_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/add_page_break", response_model=ApiResponse, tags=["Word"])
async def word_add_page_break(req: WordPageBreakRequest):
    """Add a page break to a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        doc.add_page_break()
        doc.save(req.file_path)
        return ApiResponse(message="Page break added")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/word/header_footer", response_model=ApiResponse, tags=["Word"])
async def word_header_footer(req: AddWordHeaderFooterRequest):
    """Set header and/or footer text in a Word document."""
    from docx import Document

    try:
        doc = Document(req.file_path)
        section = doc.sections[0]
        if req.header_text is not None:
            header = section.header
            header.is_linked_to_previous = False
            if header.paragraphs:
                header.paragraphs[0].text = req.header_text
            else:
                header.add_paragraph(req.header_text)
        if req.footer_text is not None:
            footer = section.footer
            footer.is_linked_to_previous = False
            if footer.paragraphs:
                footer.paragraphs[0].text = req.footer_text
            else:
                footer.add_paragraph(req.footer_text)
        doc.save(req.file_path)
        return ApiResponse(message="Header/footer updated")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  Excel endpoints
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@app.post("/excel/create", response_model=ApiResponse, tags=["Excel"])
async def excel_create(req: CreateExcelRequest):
    """Create a new Excel workbook with one or more sheets."""
    from openpyxl import Workbook

    try:
        wb = Workbook()
        default_sheet = wb.active
        for i, sheet_data in enumerate(req.sheets):
            if i == 0:
                ws = default_sheet
                ws.title = sheet_data.name
            else:
                ws = wb.create_sheet(title=sheet_data.name)
            for col, header in enumerate(sheet_data.headers, 1):
                ws.cell(row=1, column=col, value=header)
            for row_idx, row in enumerate(sheet_data.rows, 2):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
            if sheet_data.hidden:
                ws.sheet_state = "hidden"
        Path(req.file_path).parent.mkdir(parents=True, exist_ok=True)
        wb.save(req.file_path)
        return ApiResponse(message=f"Created {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/write_data", response_model=ApiResponse, tags=["Excel"])
async def excel_write_data(req: WriteExcelDataRequest):
    """Write data to a sheet in an existing workbook (creates sheet if needed)."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name in wb.sheetnames:
            ws = wb[req.sheet_name]
        else:
            ws = wb.create_sheet(title=req.sheet_name)
        for col, header in enumerate(req.headers, 1):
            ws.cell(row=1, column=col, value=header)
        for row_idx, row in enumerate(req.rows, 2):
            for col_idx, value in enumerate(row, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        if req.hidden:
            ws.sheet_state = "hidden"
        wb.save(req.file_path)
        return ApiResponse(message=f"Data written to {req.sheet_name} in {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/read", response_model=ApiResponse, tags=["Excel"])
async def excel_read(req: ReadExcelRequest):
    """Read data from an Excel workbook."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path, data_only=True)
        result = {}
        sheets_to_read = [req.sheet_name] if req.sheet_name else wb.sheetnames
        for sheet_name in sheets_to_read:
            if sheet_name not in wb.sheetnames:
                continue
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(cell) if cell is not None else "" for cell in row])
            result[sheet_name] = {"rows": rows, "hidden": ws.sheet_state == "hidden"}
        return ApiResponse(message="OK", data={"sheets": result})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/apply_formula", response_model=ApiResponse, tags=["Excel"])
async def excel_apply_formula(req: ApplyFormulaRequest):
    """Write a formula to a cell (e.g. =SUM(A1:A10))."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        ws = wb[req.sheet_name]
        ws[req.cell] = req.formula
        wb.save(req.file_path)
        return ApiResponse(message=f"Formula set in {req.cell}: {req.formula}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/format_range", response_model=ApiResponse, tags=["Excel"])
async def excel_format_range(req: FormatRangeRequest):
    """Format cells: font, fill color, number format, hide rows."""
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        ws = wb[req.sheet_name]

        cell_range = req.start_cell if not req.end_cell else f"{req.start_cell}:{req.end_cell}"
        for row in ws[cell_range]:
            cells = row if isinstance(row, tuple) else (row,)
            for cell in cells:
                if req.bold is not None or req.italic is not None or req.font_size or req.font_color:
                    font_kwargs: Dict[str, Any] = {}
                    if req.bold is not None:
                        font_kwargs["bold"] = req.bold
                    if req.italic is not None:
                        font_kwargs["italic"] = req.italic
                    if req.font_size:
                        font_kwargs["size"] = req.font_size
                    if req.font_color:
                        font_kwargs["color"] = req.font_color.lstrip("#")
                    cell.font = Font(**font_kwargs)
                if req.fill_color:
                    cell.fill = PatternFill(start_color=req.fill_color.lstrip("#"), end_color=req.fill_color.lstrip("#"), fill_type="solid")
                if req.number_format:
                    cell.number_format = req.number_format

        # Hide rows if requested
        if req.hidden:
            from openpyxl.utils import range_boundaries
            min_col, min_row, max_col, max_row = range_boundaries(cell_range)
            for r in range(min_row, max_row + 1):
                ws.row_dimensions[r].hidden = True

        wb.save(req.file_path)
        return ApiResponse(message=f"Formatted range {cell_range}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/merge_cells", response_model=ApiResponse, tags=["Excel"])
async def excel_merge_cells(req: MergeCellsRequest):
    """Merge a range of cells."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        ws = wb[req.sheet_name]
        ws.merge_cells(f"{req.start_cell}:{req.end_cell}")
        wb.save(req.file_path)
        return ApiResponse(message=f"Merged {req.start_cell}:{req.end_cell}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/delete_sheet", response_model=ApiResponse, tags=["Excel"])
async def excel_delete_sheet(req: SheetOpRequest):
    """Delete a sheet from a workbook."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        del wb[req.sheet_name]
        wb.save(req.file_path)
        return ApiResponse(message=f"Deleted sheet '{req.sheet_name}'")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/rename_sheet", response_model=ApiResponse, tags=["Excel"])
async def excel_rename_sheet(req: SheetOpRequest):
    """Rename a sheet in a workbook."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        wb[req.sheet_name].title = req.new_name
        wb.save(req.file_path)
        return ApiResponse(message=f"Renamed '{req.sheet_name}' -> '{req.new_name}'")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/insert_rows", response_model=ApiResponse, tags=["Excel"])
async def excel_insert_rows(req: InsertRowsColsRequest):
    """Insert rows at a given index."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        ws = wb[req.sheet_name]
        ws.insert_rows(req.index, req.count)
        wb.save(req.file_path)
        return ApiResponse(message=f"Inserted {req.count} row(s) at index {req.index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel/insert_cols", response_model=ApiResponse, tags=["Excel"])
async def excel_insert_cols(req: InsertRowsColsRequest):
    """Insert columns at a given index."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(req.file_path)
        if req.sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=400, detail=f"Sheet '{req.sheet_name}' not found")
        ws = wb[req.sheet_name]
        ws.insert_cols(req.index, req.count)
        wb.save(req.file_path)
        return ApiResponse(message=f"Inserted {req.count} column(s) at index {req.index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  PowerPoint endpoints
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@app.post("/pptx/create", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_create(req: CreatePptxRequest):
    """Create a new PowerPoint presentation with slides."""
    from pptx import Presentation

    try:
        prs = Presentation()
        for slide_data in req.slides:
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            if slide_data.title:
                slide.shapes.title.text = slide_data.title
            if slide_data.content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.content
            if slide_data.notes:
                slide.notes_slide.notes_text_frame.text = slide_data.notes
        Path(req.file_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(req.file_path)
        return ApiResponse(message=f"Created {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/add_slide", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_add_slide(req: AddSlideRequest):
    """Add a slide to an existing PowerPoint presentation."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if req.title:
            slide.shapes.title.text = req.title
        if req.content and len(slide.placeholders) > 1:
            slide.placeholders[1].text = req.content
        if req.notes:
            slide.notes_slide.notes_text_frame.text = req.notes
        prs.save(req.file_path)
        return ApiResponse(message=f"Slide added to {req.file_path}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/read", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_read(req: ReadPptxRequest):
    """Read all text and notes from a PowerPoint presentation."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text_frame.text)
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
            slides.append({
                "slide_number": i,
                "text": "\n".join(slide_text),
                "notes": notes,
            })
        return ApiResponse(message="OK", data={"slides": slides})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/add_table", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_add_table(req: AddPptxTableRequest):
    """Add a table to a specific slide."""
    from pptx import Presentation
    from pptx.util import Inches

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        slide = prs.slides[req.slide_index]
        num_cols = len(req.headers) if req.headers else (len(req.rows[0]) if req.rows else 1)
        num_rows = (1 if req.headers else 0) + len(req.rows)
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(req.left), Inches(req.top),
            Inches(req.width), Inches(req.height),
        )
        table = table_shape.table
        row_offset = 0
        if req.headers:
            for j, h in enumerate(req.headers):
                table.cell(0, j).text = h
            row_offset = 1
        for i, row_data in enumerate(req.rows):
            for j, val in enumerate(row_data):
                if j < num_cols:
                    table.cell(i + row_offset, j).text = str(val)
        prs.save(req.file_path)
        return ApiResponse(message=f"Table added to slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/update_slide", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_update_slide(req: UpdateSlideContentRequest):
    """Update title, content, or notes on an existing slide."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        slide = prs.slides[req.slide_index]
        if req.title is not None and slide.shapes.title:
            slide.shapes.title.text = req.title
        if req.content is not None and len(slide.placeholders) > 1:
            slide.placeholders[1].text = req.content
        if req.notes is not None:
            slide.notes_slide.notes_text_frame.text = req.notes
        prs.save(req.file_path)
        return ApiResponse(message=f"Updated slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/delete_slide", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_delete_slide(req: DeleteSlideRequest):
    """Delete a slide by index."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        rId = prs.slides._sldIdLst[req.slide_index].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[req.slide_index]
        prs.save(req.file_path)
        return ApiResponse(message=f"Deleted slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/duplicate_slide", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_duplicate_slide(req: DuplicateSlideRequest):
    """Duplicate a slide by index (appended at end)."""
    from pptx import Presentation
    from lxml import etree

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        source = prs.slides[req.slide_index]
        layout = source.slide_layout
        new_slide = prs.slides.add_slide(layout)
        # Copy all shapes from source
        for shape in source.shapes:
            el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.append(el)
        # Remove default placeholder shapes that came with layout
        # (keep only the copied ones)
        prs.save(req.file_path)
        return ApiResponse(message=f"Duplicated slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/set_notes", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_set_notes(req: SetSlideNotesRequest):
    """Set speaker notes on an existing slide."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        slide = prs.slides[req.slide_index]
        slide.notes_slide.notes_text_frame.text = req.notes
        prs.save(req.file_path)
        return ApiResponse(message=f"Notes set on slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/get_slide_info", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_get_slide_info(req: GetSlideInfoRequest):
    """Get detailed info about a specific slide."""
    from pptx import Presentation

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        slide = prs.slides[req.slide_index]
        shapes_info = []
        for shape in slide.shapes:
            info: Dict[str, Any] = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            if shape.has_text_frame:
                info["text"] = shape.text_frame.text
            if shape.has_table:
                info["table_rows"] = len(shape.table.rows)
                info["table_cols"] = len(shape.table.columns)
            shapes_info.append(info)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        return ApiResponse(
            message="OK",
            data={
                "slide_index": req.slide_index,
                "shapes": shapes_info,
                "notes": notes,
                "layout_name": slide.slide_layout.name,
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx/add_shape", response_model=ApiResponse, tags=["PowerPoint"])
async def pptx_add_shape(req: AddPptxShapeRequest):
    """Add a shape (rectangle, oval, etc.) with optional text to a slide."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    try:
        prs = Presentation(req.file_path)
        if req.slide_index < 0 or req.slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Slide index {req.slide_index} out of range")
        slide = prs.slides[req.slide_index]

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "cloud": MSO_SHAPE.CLOUD,
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
        }
        mso = shape_map.get(req.shape_type.lower(), MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            mso, Inches(req.left), Inches(req.top), Inches(req.width), Inches(req.height),
        )
        if req.text:
            shape.text_frame.text = req.text
        if req.fill_color:
            color_hex = req.fill_color.lstrip("#")
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)

        prs.save(req.file_path)
        return ApiResponse(message=f"Shape '{req.shape_type}' added to slide {req.slide_index}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── Health check ─────────────────────────────────────────────


@app.get("/health")
async def health():
    return {"status": "ok"}


# ── Main ─────────────────────────────────────────────────────

if __name__ == "__main__":
    port = 8007
    if "--port" in sys.argv:
        idx = sys.argv.index("--port")
        if idx + 1 < len(sys.argv):
            port = int(sys.argv[idx + 1])
    uvicorn.run(app, host="0.0.0.0", port=port)
